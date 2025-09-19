import csv
import re
from dataclasses import dataclass
from typing import List, Dict, Tuple, Optional

@dataclass
class QAIssue:
    uid: str
    issue_type: str
    details: str
    source: str
    target: str

    def to_dict(self):
        return {
            "ID": self.uid,
            "Issue Type": self.issue_type,
            "Details": self.details,
            "Source": self.source,
            "Target": self.target,
        }


def extract_segments(file_path: str, is_source: bool = True) -> List[str]:
    """Extract ordered segments from Word, Excel, PowerPoint, PDF (source only)."""
    segments = []

    if file_path.endswith(".xlsx"):
        import pandas as pd
        df = pd.read_excel(file_path)
        col = "source" if is_source else "target"
        if col in df.columns:
            segments = df[col].astype(str).fillna("").tolist()
        else:
            # fallback: take first column
            segments = df.iloc[:, 0].astype(str).fillna("").tolist()

    elif file_path.endswith(".docx"):
        from docx import Document
        doc = Document(file_path)
        segments = [para.text for para in doc.paragraphs if para.text.strip()]

    elif file_path.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    segments.append(shape.text)

    elif file_path.endswith(".pdf"):
        if not is_source:
            raise ValueError("PDF allowed only for source, not for target.")
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                segments.append(text)

    else:
        raise ValueError("Unsupported file format. Allowed: docx, xlsx, pptx, pdf (source only).")

    return segments


def load_pairs(src_file: str, tgt_file: Optional[str] = None) -> List[Dict[str, str]]:
    """
    Align segments from source and target files.
    Returns list of dicts {id, source, target}.
    """
    src_segments = extract_segments(src_file, is_source=True)
    tgt_segments = extract_segments(tgt_file, is_source=False) if tgt_file else []

    pairs = []
    max_len = max(len(src_segments), len(tgt_segments))

    for i in range(max_len):
        pairs.append({
            "id": str(i + 1),
            "source": src_segments[i] if i < len(src_segments) else "",
            "target": tgt_segments[i] if i < len(tgt_segments) else "",
        })

    return pairs


def load_glossary(file_path: str) -> Dict[str, str]:
    glossary = {}
    with open(file_path, encoding="utf-8-sig") as f:
        reader = csv.DictReader(f)
        for row in reader:
            src = row.get("source") or row.get("Source")
            tgt = row.get("target") or row.get("Target")
            if src and tgt:
                glossary[src.strip()] = tgt.strip()
    return glossary


# --- QA Checks (same as before) ---
def check_placeholders(uid: str, src: str, tgt: str) -> Optional[QAIssue]:
    pattern = r"%\w+|\{[^}]+\}|\$\w+"
    src_tokens = re.findall(pattern, src)
    tgt_tokens = re.findall(pattern, tgt)
    if sorted(src_tokens) != sorted(tgt_tokens):
        return QAIssue(uid, "PLACEHOLDER_MISMATCH",
                       f"Source={src_tokens}, Target={tgt_tokens}",
                       src, tgt)
    return None


def check_numbers(uid: str, src: str, tgt: str) -> Optional[QAIssue]:
    nums_src = re.findall(r"\d+", src)
    nums_tgt = re.findall(r"\d+", tgt)
    if sorted(nums_src) != sorted(nums_tgt):
        return QAIssue(uid, "NUM_MISMATCH",
                       f"Source={nums_src}, Target={nums_tgt}",
                       src, tgt)
    return None


def check_tags(uid: str, src: str, tgt: str) -> Optional[QAIssue]:
    tags_src = re.findall(r"<[^>]+>", src)
    tags_tgt = re.findall(r"<[^>]+>", tgt)
    if sorted(tags_src) != sorted(tags_tgt):
        return QAIssue(uid, "TAG_MISMATCH",
                       f"Source={tags_src}, Target={tags_tgt}",
                       src, tgt)
    return None


def check_empty(uid: str, src: str, tgt: str) -> Optional[QAIssue]:
    if src.strip() and not tgt.strip():
        return QAIssue(uid, "EMPTY_TARGET", "Target is empty", src, tgt)
    return None


def check_glossary(uid: str, src: str, tgt: str, glossary: Dict[str, str]) -> Optional[QAIssue]:
    for s_term, t_term in glossary.items():
        if s_term in src and t_term not in tgt:
            return QAIssue(uid, "GLOSSARY_MISMATCH",
                           f"Expected '{s_term}' -> '{t_term}'",
                           src, tgt)
    return None


def run_checks(pairs: List[Dict[str, str]],
               glossary: Optional[Dict[str, str]] = None,
               config: Optional[Dict] = None,
               length_ratio_limits: Tuple[float, float] = (0.5, 3.0)) -> Tuple[List[QAIssue], Dict]:
    issues: List[QAIssue] = []
    stats = {"total": len(pairs), "issues": 0}

    for p in pairs:
        uid, src, tgt = p.get("id"), p.get("source", ""), p.get("target", "")

        # Length ratio check
        if src and tgt:
            ratio = len(tgt) / max(1, len(src))
            if not (length_ratio_limits[0] <= ratio <= length_ratio_limits[1]):
                issues.append(QAIssue(uid, "LENGTH_RATIO",
                                      f"Ratio={ratio:.2f}", src, tgt))

        for check in [check_placeholders, check_numbers, check_tags, check_empty]:
            issue = check(uid, src, tgt)
            if issue:
                issues.append(issue)

        if glossary:
            issue = check_glossary(uid, src, tgt, glossary)
            if issue:
                issues.append(issue)

    stats["issues"] = len(issues)
    return issues, stats
